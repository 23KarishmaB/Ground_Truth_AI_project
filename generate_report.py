import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from dotenv import load_dotenv
import openai

BASE_DIR = Path(__file__).resolve().parents[1]
OUTPUT_DIR = BASE_DIR / "output"
OUTPUT_DIR.mkdir(exist_ok=True)

load_dotenv()
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
if OPENAI_API_KEY:
    openai.api_key = OPENAI_API_KEY


def call_ai(prompt: str) -> str:
    if not OPENAI_API_KEY:
        return "AI key not set. To enable AI summaries, add OPENAI_API_KEY in a .env file."

    try:
        resp = openai.ChatCompletion.create(
            model="gpt-4o-mini",
            messages=[{"role": "system", "content": "You are an analytics assistant."},
                      {"role": "user", "content": prompt}],
            max_tokens=200,
            temperature=0.2
        )
        return resp["choices"][0]["message"]["content"].strip()
    except Exception as e:
        return f"AI call failed: {e}"


def build_report(metrics: dict, imgs: dict) -> Path:
    prs = Presentation()
    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "Automated Insight Engine — Spotify Analysis"
    try:
        title_slide.placeholders[1].text = "Auto-generated report"
    except Exception:
        pass

    # Summary slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Dataset Summary"
    tf = slide.shapes.placeholders[1].text_frame
    for k, v in metrics.items():
        p = tf.add_paragraph()
        p.text = f"{k}: {v}"

    # Add images
    for name, path in imgs.items():
        if path:
            s = prs.slides.add_slide(prs.slide_layouts[5])
            s.shapes.title.text = name.replace("_", " ").title()
            s.shapes.add_picture(path, Inches(0.5), Inches(1.2), width=Inches(9))
            # Optionally add AI insight below each image
            if OPENAI_API_KEY:
                insight = call_ai(f"Provide 2 short insights for the chart titled '{name}'.")
                tx = s.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(9), Inches(1.0))
                tx.text_frame.text = "AI Insight: " + insight

    # AI insights slide
    if OPENAI_API_KEY:
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = "Executive AI Insights"
        text = call_ai("Provide 5 high-level executive insights and next steps for a music streaming product using popularity and feature trends.")
        s.shapes.placeholders[1].text = text

    out = OUTPUT_DIR / "report.pptx"
    prs.save(out)
    print(f"[✔] Report saved to {out}")
    return out
